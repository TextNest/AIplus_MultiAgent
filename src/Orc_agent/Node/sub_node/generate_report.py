import os
import io
import time
import pandas as pd
import base64

from pptx import Presentation
from pptx.util import Inches, Pt
from xhtml2pdf import pisa
import docx

from ...core.llm_factory import LLMFactory
from ...core.observe import langfuse_session, merge_runnable_config, observe
from langchain_core.runnables import RunnableConfig
from src.Orc_agent.core.prompts import (
    REPORT_PROMPT_GENERAL, 
    REPORT_PROMPT_DECISION, 
    REPORT_PROMPT_MARKETING,
    HTML_WRAPPER_PROMPT,
    REPORT_STYLE_CLASSIFICATION_PROMPT,
    OVERALL_PPT_PROMPT,
    CHART_PPT_PROMPT
)
from ...State.state import ReportState

from src.Orc_agent.core.logger import logger
from typing import Literal, List

@observe(name="Supervisor")
def report_supervisor(state: ReportState) -> ReportState:
    """
    Supervisor node that decides the next step in report generation.
    """
    final_report = state.get("final_report")
    report_format = state.get("report_format", ["html"])
    report_style = state.get("report_style")
    logger.info(f"현재 등록된 보고서 형식 : {report_format}, 스타일: {report_style}")

    
    # Ensure report_format is a list of lowercase strings
    if isinstance(report_format, str):
        report_format = [report_format]
    report_format = [f.lower() for f in report_format]
    
    generated_formats = state.get("generated_formats", [])
    
    # helper to safely check format validity
    def needs_format(fmt):
        return fmt in report_format and fmt not in generated_formats

    # 1. If HTML content is missing, generate it first
    # 먼저 HTML 변환 계열(pdf, html, docs)이 하나라도 요청 목록에 있는지 확인
    needs_html_family = any(fmt in ["pdf", "html", "docx"] for fmt in report_format)
    # 필요한데 아직 안 만들어졌다면 -> generate_content로 보냄
    if needs_html_family and not final_report:
        # 보고서 스타일 분류
        # If report style is not decided or set to "AI 자동 판단 (추천)", classify it first
        if not report_style or report_style == "AI 자동 판단 (추천)":
            return {"next_worker": "classify_report_style"}
        return {"next_worker": "generate_content"}
    
    # 2. Check for requested formats that haven't been generated yet
    # html이 준비되었거나, 애초에 PPTX만 요청한 경우 순서대로 아직 안 만든 파일들 생성
    if "pdf" in report_format and "pdf" not in generated_formats:
        return {"next_worker": "create_pdf", "generated_formats": ["pdf"]}
        
    if "docx" in report_format and "docx" not in generated_formats:
        return {"next_worker": "create_docx", "generated_formats": ["docx"]}
        
    if "pptx" in report_format and "pptx" not in generated_formats:
        return {"next_worker": "create_pptx", "generated_formats": ["pptx"]}
        
    # 3. If all done (or just html requested), finish
    return {"next_worker": "FINISH"}

@observe(name="classify_report_style")
def classify_report_style(state: ReportState, config: RunnableConfig) -> ReportState:
    """
    Classifies the analysis results to choose the best report style using overall_insight.
    """
    analysis_results = state.get("analysis_results", {})
    
    if not analysis_results:
        return {"report_style": "일반 리포트", "steps_log": ["[Report] No analysis results, defaulting to 일반 리포트"]}
    
    try:
        node_conf = state.get("node_models", {}).get("report_style_node", {})
        provider = node_conf.get("provider") or "google"
        model_name = node_conf.get("model") or "gemini-2.5-flash"
        
        llm, callbacks = LLMFactory.create(
            provider=provider,
            model=model_name,
            temperature=0,
        )
        
        # Extract overall insights for classification
        overall_insights = []
        if isinstance(analysis_results, dict):
            for key, value in analysis_results.items():
                if "overall" in key.lower():
                    insight = value.get("insight", "") if isinstance(value, dict) else str(value)
                    if insight:
                        overall_insights.append(insight)
        
        all_overall_text = "\n".join(overall_insights)
        if not all_overall_text:
             # Fallback: use first few results if overall is missing
             all_overall_text = str(list(analysis_results.values())[0])[:1000]

        prompt = REPORT_STYLE_CLASSIFICATION_PROMPT.format(overall_insight=all_overall_text)

        response = llm.invoke(prompt, config=config)
        selected_style = response.content.strip()
        
        # Validation
        valid_styles = ["일반 리포트", "의사 결정 리포트", "마케팅 예산 분배 리포트"]
        if selected_style not in valid_styles:
            for style in valid_styles:
                if style in selected_style:
                    selected_style = style
                    break
            else:
                selected_style = "일반 리포트"

        logger.info(f"AI가 선택한 보고서 스타일: {selected_style}")
        return {
            "report_style": selected_style,
            "steps_log": [f"[Report] AI classified report style as: {selected_style}"]
        }
    except Exception as e:
        logger.error(f"보고서 스타일 분류 중 오류: {e}")
        return {
            "report_style": "일반 리포트",
            "steps_log": [f"[Report] Classification failed, defaulting to 일반 리포트: {str(e)}"]
        }

@observe(name="generate_content")
def generate_content(state: ReportState, config: RunnableConfig) -> ReportState:
    """
    Generates report content in HTML format using LLM.
    """
    analysis_results = state.get("analysis_results", {})
    file_path = state.get("file_path", "Data")
    figure_list = state.get("figure_list", [])
    
    if not analysis_results:
        return {
            "final_report": "# Error\n\nNo analysis results available.",
            "steps_log": ["[Report] ERROR: No analysis results"]
        }

    try:
        # LLM Setup
        node_conf = state.get("node_models", {}).get("report_gen_node", {})
        provider = node_conf.get("provider") or "google"
        model_name = node_conf.get("model") or "gemini-2.5-flash"
        
        llm, callbacks = LLMFactory.create(
            provider=provider,
            model=model_name,
            temperature=0.3,
        )
        
        # Data Context
        data_summary = state.get("formatted_output", "")
        
        # Visualization Context
        figure = ""
        base64_mapping = {}  # 플레이스홀더 맵핑용
        if figure_list:
            figure = "### Key Visualizations\n"
            for i, fig in enumerate(figure_list):
                if os.path.exists(fig):
                    # 이미지를 열어서 base64로 인코딩
                    with open(fig, "rb") as image_file:
                        encoded_string = base64.b64encode(image_file.read()).decode()
                    placeholder = f"[FIGURE_{i}_PLACEHOLDER]"
                    base64_mapping[placeholder] = f"data:image/png;base64,{encoded_string}"
                    # 로컬 src 경로 대신 매직 스트링(data:image/png...) 사용
                    figure += f"<img src='{placeholder}' alt='{os.path.basename(fig)}' style='max-width: 100%;'>\n"
        
        # Process Analysis Results
        processed_results = []
        if isinstance(analysis_results, dict):
            for key, value in analysis_results.items():
                insight = ""
                if isinstance(value, dict):
                    insight = value.get("insight", "")
                else:
                    insight = str(value)
                
                if insight:
                    if "overall" in key.lower():
                        processed_results.append(f"### [Overall Insight]\n{insight}")
                    else:
                        processed_results.append(f"### [Insight for {key}]\n{insight}")
        elif isinstance(analysis_results, list):
            for res in analysis_results:
                processed_results.append(str(res))
        
        all_results = "\n\n---\n\n".join(processed_results)
        
        prompt_map = {
            "일반 리포트": REPORT_PROMPT_GENERAL,
            "의사 결정 리포트": REPORT_PROMPT_DECISION,
            "마케팅 예산 분배 리포트": REPORT_PROMPT_MARKETING
        }
        selected_style = state.get("report_style", "일반 리포트")
        prompt_template = prompt_map.get(selected_style, REPORT_PROMPT_GENERAL)
        # 1단계: 핵심 지시문
        core_instructions = prompt_template.format(
            data_summary=data_summary,
            all_results=all_results,
            figure=figure
        )

        # 2단계: HTML 껍데기에 핵심 지시문 넣기
        prompt = HTML_WRAPPER_PROMPT.format(
            core_instructions=core_instructions,
            selected_style=selected_style
        )

        # LLM에게 완성된 프롬프트 전달
        with langfuse_session(session_id="generate-report", tags=["generate_report"]) as lf_metadata:
            invoke_cfg = merge_runnable_config(
                config,
                callbacks=callbacks,
                metadata=lf_metadata,
            )
            time.sleep(2) # 할당량 한도(Rate Limit) 방어
            response = llm.invoke(prompt, config=invoke_cfg)
        
        # Handle DummyLLM response
        if hasattr(response, 'content'):
            content = response.content
        else:
            content = str(response)

        if isinstance(content, list):
            # LLM 멀티모달 응답의 경우 dict의 'text' 추출
            parts = []
            for part in content:
                if isinstance(part, dict) and "text" in part:
                    parts.append(part["text"])
                else:
                    parts.append(str(part))
            content = "".join(parts)
            
        # 1. 마크다운 기호 제거 및 깔끔한 정리 (HTML 태그만 남기기)
        import re
        content = re.sub(r"^```[a-zA-Z]*\n", "", content.strip(), flags=re.IGNORECASE)
        content = re.sub(r"\n```$", "", content)
        content = content.replace("```html", "").replace("```", "")
        
        # 1-1. LLM의 출력물이 텍스트 문자열(literal '\\n')을 포함하는 경우 실제 개행으로 변경
        content = content.replace("\\n", "\n")
        
        # 2. 플레이스홀더를 실제 Base64 데이터로 교체
        for placeholder, base64_data in base64_mapping.items():
            # 정상적인 태그 교체
            content = content.replace(placeholder, base64_data)
            # LLM이 특수문자를 URL 인코딩(%5B, %5D)했을 경우 교체
            url_encoded = placeholder.replace("[", "%5B").replace("]", "%5D")
            content = content.replace(url_encoded, base64_data)
            # 대소문자 오류 방어
            content = content.replace(placeholder.lower(), base64_data)

        return {
            "final_report": content,
            "steps_log": ["[Report] Generated content via LLM"]
        }

    except Exception as e:
        return {
            "final_report": f"# Error\n\nReport generation failed: {str(e)}",
            "steps_log": [f"[Report] ERROR: {str(e)}"]
        }
    
@observe(name="create_pdf")
def create_pdf(state: ReportState, config: RunnableConfig) -> ReportState:
    """
    Converts HTML report to PDF.
    """
    thread_id = config["configurable"].get("thread_id", "default")
    output_dir = os.path.join("output", thread_id)
    os.makedirs(output_dir, exist_ok=True) # 폴더가 없으면 알아서 생성!
    content = state.get("final_report", "")
    if not content:
        return {"steps_log": ["[Report] PDF Generation Skipped (No Content)"]}
        
    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        
        # 1. Font Source (Cross-platform)
        import platform
        font_name = "Helvetica"
        system_font_path = ""
        
        system = platform.system()
        if system == "Windows":
            font_name = "MalgunGothic"
            system_font_path = r"C:/Windows/Fonts/malgun.ttf"
        
        # 리눅스, Mac 또는 지정 폰트가 없을 시: matplotlib 활용 탐색
        if not os.path.exists(system_font_path):
            try:
                from matplotlib import font_manager
                font_list = font_manager.findSystemFonts(fontpaths=None, fontext='ttf')
                preferred = ['nanumgothic', 'nanumbarun', 'malgun', 'notosanscjk', 'd2coding', 'applegothic']
                for pref in preferred:
                    for f in font_list:
                        if pref in f.lower():
                            system_font_path = f
                            font_name = "KoreanFont"
                            break
                    if system_font_path and os.path.exists(system_font_path):
                        break
            except Exception:
                pass
        
        font_registered = False
        if system_font_path and os.path.exists(system_font_path):
            try:
                pdfmetrics.registerFont(TTFont(font_name, system_font_path))
                font_registered = True
            except Exception as e:
                print(f"Font registration warning: {e}")
        
        if not font_registered:
             # Fallback to a standard font if registration fails
             font_name = "Helvetica" # standard PDF font
             system_font_path = ""
             print("Using fallback font: Helvetica")

        font_face_css = ""
        if system_font_path:
            # 윈도우 역슬래시 이슈 방지를 위해 주소 포맷 변경
            clean_path = system_font_path.replace('\\', '/')
            font_face_css = f'''
                @font-face {{
                    font-family: '{font_name}';
                    src: url('{clean_path}');
                }}'''

        styled_html = f"""
        <html>
        <head>
            <meta charset="utf-8">
            <style>
                {font_face_css}
                body {{
                    font-family: '{font_name}', sans-serif;
                }}
                img {{
                    max-width: 100%;
                }}
            </style>
        </head>
        <body>
            {content}
        </body>
        </html>
        """
        
        output_path = os.path.join(output_dir, "report.pdf")
        with open(output_path, "wb") as f:
            pisa_status = pisa.CreatePDF(styled_html, dest=f, encoding='utf-8')
        
        if pisa_status.err:
            raise Exception("PDF generation failed")
            
        return {
            "steps_log": [f"[Report] Generated PDF report at {output_path}"],
            "generated_formats":["pdf"]
        }
    except Exception as e:
        return {"steps_log": [f"[Report] PDF Generation Error: {str(e)}"]}
@observe(name="create_docx")
def create_docx(state: ReportState, config: RunnableConfig) -> ReportState:
    """
    Converts HTML report to Word Document (.docx).
    """
    thread_id = config["configurable"].get("thread_id", "default")
    output_dir = os.path.join("output", thread_id)
    os.makedirs(output_dir, exist_ok=True) # 폴더가 없으면 알아서 생성!
    content_text = state.get("final_report", "")
    if not content_text:
        return {"steps_log": ["[Report] DOCX Generation Skipped (No Content)"]}
        
    try:
        import docx
        from bs4 import BeautifulSoup # HTML 파싱을 위해 필요할 수 있습니다.
        
        doc = docx.Document()
        doc.add_heading('Data Analysis Report', 0)
        
        # 1. HTML에서 텍스트 추출 후 문단 추가
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(content_text, 'html.parser')
            text_blocks = soup.get_text(separator='\n\n', strip=True)
            for para in text_blocks.split('\n\n'):
                if para.strip():
                    doc.add_paragraph(para.strip())
        except Exception as e:
            doc.add_paragraph("텍스트 파싱 오류 발생: " + str(e))
        
        # 2. 이미지(figure_list)를 순회하며 워드파일 안에 박아 넣기
        figure_list = state.get("figure_list", [])
        doc.add_heading('Visualizations', level=1)
        for fig_path in figure_list:
            if os.path.exists(fig_path):
                doc.add_picture(fig_path, width=docx.shared.Inches(5.0))
                doc.add_paragraph("\n") # 여백 느낌
        
        output_path = os.path.join(output_dir, "report.docx")
        doc.save(output_path)
            
        return {
            "steps_log": [f"[Report] Generated Word report at {output_path}"],
            "generated_formats":["docx"]
        }
    except Exception as e:
        return {"steps_log": [f"[Report] DOCX Generation Error: {str(e)}"]}
    
@observe(name="create_pptx")
def create_pptx(state: ReportState, config: RunnableConfig) -> ReportState:
    """
    Generates PowerPoint report.
    """
    thread_id = config["configurable"].get("thread_id", "default")
    output_dir = os.path.join("output", thread_id)
    os.makedirs(output_dir, exist_ok=True) # 폴더가 없으면 알아서 생성!
    analysis_results = state.get("analysis_results", {})
    figure_list = state.get("figure_list", [])
    
    try:
        # LLM Setup
        node_conf = state.get("node_models", {}).get("report_gen_node", {})
        provider = node_conf.get("provider") or "google"
        model_name = node_conf.get("model") or "gemini-2.5-flash"
        
        llm, callbacks = LLMFactory.create(
            provider=provider,
            model=model_name,
            temperature=0.3,
        )

        prs = Presentation()
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Data Analysis Report"
        subtitle.text = "Generated by AIplus MultiAgent"
        
        # Content Slides
        if analysis_results:
            overall_text_raw = ""
            chart_insights = {}  # { "figure_0_0.png": "인사이트 전문", ... }
            for key, value_dict in analysis_results.items():
                if isinstance(value_dict, dict):
                    insight_text = value_dict.get("insight", "")
                else:
                    insight_text = str(value_dict)
                
                # overall로 시작하는 키는 전체 요약용으로
                if key.startswith("overall"):
                    overall_text_raw += insight_text + "\n"
                # 그 외(.png 등)는 개별 차트용으로 분류
                else:
                    chart_insights[key] = insight_text

        if overall_text_raw:
            # LLM 호출 전에 약간 대기
            time.sleep(2)
            prompt = OVERALL_PPT_PROMPT.format(text=overall_text_raw)
            response = llm.invoke(prompt)
            content_val = response.content
            if isinstance(content_val, list):
                content_val = "".join([c.get("text", "") if isinstance(c, dict) else str(c) for c in content_val])
            summary_bullet_points = content_val.strip()
            
            # 슬라이드 생성
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = "Executive Summary"
            tf = slide.placeholders[1].text_frame
            tf.text = summary_bullet_points

        # Figure Slides
        for fig_path in figure_list:
            if os.path.exists(fig_path):
                # 1. 이미지 파일명 추출 (예: figure_0_0.png)
                img_filename = os.path.basename(fig_path)
                
                # 2. 딕셔너리에서 해당 차트의 원본 인사이트 가져오기
                raw_chart_insight = chart_insights.get(img_filename, "")
                
                # 3. 차트 요약 생성
                chart_bullet_points = "요약 정보 없음" # 기본값
                if raw_chart_insight:
                    time.sleep(3) # 버스트 호출 방지용 안전 지연
                    prompt = CHART_PPT_PROMPT.format(text=raw_chart_insight)
                    response = llm.invoke(prompt)
                    content_val = response.content
                    if isinstance(content_val, list):
                        content_val = "".join([c.get("text", "") if isinstance(c, dict) else str(c) for c in content_val])
                    chart_bullet_points = content_val.strip()
                # 4. 슬라이드 레이아웃 설정 (텍스트 + 컨텐츠 레이아웃 추천: 레이아웃 [8])
                # 또는 빈 레이아웃[6]에 텍스트박스와 그림을 직접 좌표 지정해서 넣기 (직접 지정 예시)
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # (왼쪽) 요약 텍스트 추가
                # left=0.5, top=1.0, width=4.0, height=5.5 (인치)
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4), Inches(5.5))
                tf = txBox.text_frame
                tf.word_wrap = True # 텍스트 안 잘리게 줄바꿈 강제
                tf.text = f"{img_filename} Analysis:\n\n{chart_bullet_points}"
                
                # (오른쪽) 이미지 추가 (비율이 안 깨지도록 위치 조정)
                # left=4.5, top=1.0, height=5.0 (가로 비율은 자동)
                slide.shapes.add_picture(fig_path, Inches(4.5), Inches(1), height=Inches(5))
        
        output_path = os.path.join(output_dir, "report.pptx")
        prs.save(output_path)
        
        return {
            "steps_log": [f"[Report] Generated PowerPoint report at {output_path}"],
            "generated_formats":["pptx"]
        }
    except Exception as e:
        print(f"🔥 PPTX 생성 진짜 에러 원인: {e}") 
        return {"steps_log": [f"[Report] PPTX Generation Error: {str(e)}"]}
