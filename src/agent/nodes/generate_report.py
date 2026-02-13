"""
Generate Report Node
Roles: Analyzer, Reporter
Input: analysis_results, clean_data, figure_list
Output: final_report, steps_log
"""
import os
import io
import pandas as pd
import markdown

from pptx import Presentation
from pptx.util import Inches, Pt
from xhtml2pdf import pisa

from ..state import AgentState
from ..prompt_engineering.prompts import REPORT_PROMPT
from ...core.llm_factory import LLMFactory
from ...core.observe import langfuse_session


def generate_report(state: AgentState) -> AgentState:
    """
    Generates a report based on analysis results.
    Supports Markdown, HTML, PDF, and PowerPoint formats.
    """
    analysis_results = state.get("analysis_results", [])
    clean_data = state.get("clean_data")
    file_path = state.get("file_path", "Data")
    figure_list = state.get("figure_list", [])
    report_format = state.get("report_format", "markdown").lower()
    
    if not analysis_results:
        return {
            "final_report": "# Error\n\nNo analysis results available.",
            "steps_log": ["[Report] ERROR: No analysis results"]
        }
    
    try:  
        # Step 1: Generate Markdown Content via LLM
        markdown_content = _generate_markdown_content(analysis_results, clean_data, file_path, figure_list)
        
        # Step 2: Convert to Requested Format
        final_output = markdown_content # Default
        log_message = "[Report] Generated Markdown report successfully"
        
        output_dir = "output"
        os.makedirs(output_dir, exist_ok=True)
        
        if report_format == "html":
            html_content = markdown.markdown(markdown_content)
            # Add basic styling
            styled_html = f"<html><body><style>body {{ font-family: sans-serif; max-width: 800px; margin: auto; padding: 20px; }} img {{ max-width: 100%; }}</style>{html_content}</body></html>"
            
            output_path = os.path.join(output_dir, "report.html")
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(styled_html)
            
            final_output = f"Report saved to: {output_path}"
            log_message = f"[Report] Generated HTML report at {output_path}"
            
        elif report_format == "pdf":
            html_content = markdown.markdown(markdown_content)
            # Add basic styling for PDF
            styled_html = f"<html><body><style>body {{ font-family: sans-serif; }} img {{ max-width: 100%; }}</style>{html_content}</body></html>"
            
            output_path = os.path.join(output_dir, "report.pdf")
            with open(output_path, "wb") as f:
                pisa_status = pisa.CreatePDF(html_content, dest=f)
            
            if pisa_status.err:
                raise Exception("PDF generation failed")
                
            final_output = f"Report saved to: {output_path}"
            log_message = f"[Report] Generated PDF report at {output_path}"
            
        elif report_format == "pptx":
            output_path = os.path.join(output_dir, "report.pptx")
            _generate_pptx(analysis_results, figure_list, output_path)
            
            final_output = f"Report saved to: {output_path}"
            log_message = f"[Report] Generated PowerPoint report at {output_path}"

        return {
            "final_report": final_output,
            "steps_log": [log_message]
        }
    except Exception as e:
        import traceback
        traceback.print_exc()
        return {
            "final_report": f"# Error\n\nReport generation failed: {str(e)}",
            "steps_log": [f"[Report] ERROR: {str(e)}"]
        }

def _generate_markdown_content(results, data, file_path, figures):
    # LLM Setup
    llm, callbacks = LLMFactory.create(
        provider="google",
        model="gemini-2.5-flash",
        temperature=0.3,
    )
    
    # Data Context
    data_summary = ""
    if data:
        df = pd.DataFrame(data)
        data_summary = f"""
- Data Source: {file_path}
- Rows: {len(df):,}
- Columns: {len(df.columns)}
- Column List: {', '.join(df.columns)}
"""
    
    # Visualization Context
    figure_markdown = ""
    if figures:
        figure_markdown = "### Key Visualizations\n"
        for fig in figures:
            figure_markdown += f"![{fig}]({fig})\n"
    
    all_results = "\n\n---\n\n".join(results)
    
    prompt = REPORT_PROMPT.format(
        data_summary=data_summary,
        all_results=all_results,
        figure_markdown=figure_markdown
    )

    with langfuse_session(session_id="generate-report", tags=["generate_report"]):
        response = llm.invoke(prompt, config={"callbacks": callbacks})
    
    content = response.content
    if isinstance(content, list):
        content = "".join([str(part) for part in content])
        
    # Appending figure markdown strictly might duplicate if LLM already includes it, 
    # but ensuring it's there is safer if expected. 
    # The prompt actually asks to include it, so we rely on LLM.
    # But for dummy mode or consistency, let's just return content.
    return content

def _generate_pptx(analysis_results, figure_list, output_path):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Data Analysis Report"
    subtitle.text = "Generated by AIplus MultiAgent"
    
    # Content Slides from Analysis Results
    # We'll take the first analysis result as the main content for simplicity
    if analysis_results:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Analysis Insights"
        tf = body_shape.text_frame
        
        # Naively truncate or use first chunk
        text_content = analysis_results[0].replace("```python", "").replace("```", "")
        # Limit text length roughly
        tf.text = text_content[:500] + "..." if len(text_content) > 500 else text_content

    # Figure Slides
    for fig_path in figure_list:
        if os.path.exists(fig_path):
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            left = Inches(1)
            top = Inches(1)
            height = Inches(5.5)
            slide.shapes.add_picture(fig_path, left, top, height=height)
            
    prs.save(output_path)
