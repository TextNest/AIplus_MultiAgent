
import os
import io
import pandas as pd
import markdown

from pptx import Presentation
from pptx.util import Inches, Pt
from xhtml2pdf import pisa

from src.core.llm_factory import LLMFactory, langfuse_session
from src.agent.prompt_engineering.prompts import REPORT_PROMPT
from .state import ReportState

# Ensure output directory exists
OUTPUT_DIR = "output"
os.makedirs(OUTPUT_DIR, exist_ok=True)

def generate_content(state: ReportState) -> dict:
    """
    Generates report content in Markdown format using LLM.
    """
    analysis_results = state.get("analysis_results", [])
    clean_data = state.get("clean_data")
    file_path = state.get("file_path", "Data")
    figure_list = state.get("figure_list", [])
    
    if not analysis_results:
        return {
            "final_report": "# Error\n\nNo analysis results available.",
            "steps_log": ["[Report] ERROR: No analysis results"]
        }

    try:
        # LLM Setup
        llm, callbacks = LLMFactory.create(
            provider="google",
            model="gemini-2.5-flash",
            temperature=0.3,
        )
        
        # Data Context
        data_summary = ""
        if clean_data:
            df = pd.DataFrame(clean_data)
            data_summary = f"""
- Data Source: {file_path}
- Rows: {len(df):,}
- Columns: {len(df.columns)}
- Column List: {', '.join(df.columns)}
"""
        
        # Visualization Context
        figure_markdown = ""
        if figure_list:
            figure_markdown = "### Key Visualizations\n"
            for fig in figure_list:
                figure_markdown += f"![{fig}]({fig})\n"
        
        all_results = "\n\n---\n\n".join(analysis_results)
        
        prompt = REPORT_PROMPT.format(
            data_summary=data_summary,
            all_results=all_results,
            figure_markdown=figure_markdown
        )

        with langfuse_session(session_id="generate-report", tags=["generate_report"]):
            response = llm.invoke(prompt, config={"callbacks": callbacks})
        
        # Handle DummyLLM response
        if hasattr(response, 'content'):
            content = response.content
        else:
            content = str(response)

        if isinstance(content, list):
            content = "".join([str(part) for part in content])
            
        return {
            "final_report": content,
            "steps_log": ["[Report] Generated Markdown content via LLM"]
        }

    except Exception as e:
        return {
            "final_report": f"# Error\n\nReport generation failed: {str(e)}",
            "steps_log": [f"[Report] ERROR: {str(e)}"]
        }

def create_pdf(state: ReportState) -> dict:
    """
    Converts Markdown report to PDF.
    """
    markdown_content = state.get("final_report", "")
    if not markdown_content:
        return {"steps_log": ["[Report] PDF Generation Skipped (No Content)"]}
        
    try:
        html_content = markdown.markdown(markdown_content)
        # Add basic styling for PDF
        styled_html = f"<html><body><style>body {{ font-family: sans-serif; }} img {{ max-width: 100%; }}</style>{html_content}</body></html>"
        
        output_path = os.path.join(OUTPUT_DIR, "report.pdf")
        with open(output_path, "wb") as f:
            pisa_status = pisa.CreatePDF(html_content, dest=f)
        
        if pisa_status.err:
            raise Exception("PDF generation failed")
            
        return {
            "steps_log": [f"[Report] Generated PDF report at {output_path}"]
        }
    except Exception as e:
        return {"steps_log": [f"[Report] PDF Generation Error: {str(e)}"]}

def create_html(state: ReportState) -> dict:
    """
    Converts Markdown report to HTML.
    """
    markdown_content = state.get("final_report", "")
    if not markdown_content:
        return {"steps_log": ["[Report] HTML Generation Skipped (No Content)"]}
        
    try:
        html_content = markdown.markdown(markdown_content)
        # Add basic styling
        styled_html = f"<html><body><style>body {{ font-family: sans-serif; max-width: 800px; margin: auto; padding: 20px; }} img {{ max-width: 100%; }}</style>{html_content}</body></html>"
        
        output_path = os.path.join(OUTPUT_DIR, "report.html")
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(styled_html)
            
        return {
            "steps_log": [f"[Report] Generated HTML report at {output_path}"]
        }
    except Exception as e:
        return {"steps_log": [f"[Report] HTML Generation Error: {str(e)}"]}

def create_pptx(state: ReportState) -> dict:
    """
    Generates PowerPoint report.
    """
    analysis_results = state.get("analysis_results", [])
    figure_list = state.get("figure_list", [])
    
    try:
        prs = Presentation()
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Data Analysis Report"
        subtitle.text = "Generated by AIplus MultiAgent"
        
        # Content Slides
        if analysis_results:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = "Analysis Insights"
            tf = body_shape.text_frame
            
            text_content = analysis_results[0].replace("```python", "").replace("```", "")
            tf.text = text_content[:500] + "..." if len(text_content) > 500 else text_content

        # Figure Slides
        for fig_path in figure_list:
            if os.path.exists(fig_path):
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                left = Inches(1)
                top = Inches(1)
                height = Inches(5.5)
                slide.shapes.add_picture(fig_path, left, top, height=height)
        
        output_path = os.path.join(OUTPUT_DIR, "report.pptx")
        prs.save(output_path)
        
        return {
            "steps_log": [f"[Report] Generated PowerPoint report at {output_path}"]
        }
    except Exception as e:
        return {"steps_log": [f"[Report] PPTX Generation Error: {str(e)}"]}